from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize a new Presentation
prs = Presentation()

# Function to add a title slide with Lord of the Rings style
def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Style the title text
    title_placeholder.text_frame.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(60)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 223, 0)  # Gold color

    # Style the subtitle text
    subtitle_placeholder.text_frame.text = subtitle
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color

# Function to add content slides with Lord of the Rings style
def add_content_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(48)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 223, 0)  # Gold color

    # Set content
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content
    content_placeholder.text_frame.paragraphs[0].font.size = Pt(28)
    content_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color

# Adding the slides

# Title Slide
add_title_slide(
    "Netzwerkeinrichtung für die Arztpraxis",
    "Beschaffungsplan für Hardware und Software\nPräsentiert von: [Ihr Name]\nDatum: [Präsentationsdatum]"
)

# Slide 2: Projektübersicht
add_content_slide(
    "Projektübersicht",
    "Ziel: Aufbau eines zuverlässigen und sicheren Netzwerks für die Arztpraxis.\n"
    "Umfang:\n- 1 Server und 7 Arbeitsstationen.\n- Internetzugang, Datenaustausch und mehrere Softwarepakete.\n"
    "Budget: €20.000 (netto)."
)

# Slide 3: Netzwerkanforderungen
add_content_slide(
    "Netzwerkanforderungen",
    "Internetverbindung: DSL-Anschluss für alle Arbeitsstationen.\n"
    "Datenaustausch: Arbeitsstationen müssen nahtlos miteinander kommunizieren und Daten austauschen können.\n"
    "Software-Nutzung: Unterstützung für mehrere medizinische und Bürosoftwarepakete."
)

# Slide 4: Hardwareübersicht
add_content_slide(
    "Hardwareübersicht",
    "Server: Dell PowerEdge T40, geeignet für die Verwaltung medizinischer Daten und Software.\n"
    "Arbeitsstationen: 7 x HP ProDesk 400 G6, ausgestattet für den Betrieb der benötigten Software.\n"
    "Netzwerkausrüstung: DSL-Router (Fritz!Box 7590), Netzwerkswitch und Verkabelung.\n"
    "Peripheriegeräte: Monitore, Tastatur/Maus-Sets und ein Multifunktionsdrucker/Scanner."
)

# Slide 5: Softwareübersicht
add_content_slide(
    "Softwareübersicht",
    "Betriebssysteme: Windows Server 2019 für den Server, Windows 10 Pro für die Arbeitsstationen.\n"
    "Medizinische Software: Praxis EMR und Praxisverwaltungssoftware.\n"
    "Office Suite: Microsoft Office 365 für tägliche Aufgaben.\n"
    "Sicherheitssoftware: Antivirus und Firewall für den Datenschutz."
)

# Slide 6: Implementierungsplan
add_content_slide(
    "Implementierungsplan",
    "Netzwerkeinrichtung: Installation und Konfiguration des Servers, Verbindung der Arbeitsstationen, und Einrichtung des Routers.\n"
    "Softwareinstallation: Installation der erforderlichen Software und Konfiguration der Sicherheitseinstellungen.\n"
    "Testen: Überprüfung der Netzwerkfunktionalität und Softwareleistung.\n"
    "Schulung: IT-Grundlagenschulung für das Personal zur effizienten Nutzung des neuen Systems."
)

# Slide 7: Budgetaufteilung
add_content_slide(
    "Budgetaufteilung",
    "Hardware: €9.240\n"
    "Software: €7.750\n"
    "Sonstiges (Verkabelung, Peripheriegeräte): €2.000\n"
    "Gesamt: €18.990 (innerhalb des Budgets von €20.000)"
)

# Slide 8: Vorteile des neuen Systems
add_content_slide(
    "Vorteile des neuen Systems",
    "Verbesserte Effizienz: Schnelleres und zuverlässigeres System.\n"
    "Erhöhte Sicherheit: Schutz sensibler Patientendaten.\n"
    "Skalierbarkeit: System kann bei Bedarf erweitert werden."
)

# Slide 9: Fazit
add_content_slide(
    "Fazit",
    "Zusammenfassung der wichtigsten Punkte.\n"
    "Einladung zu Fragen."
)

# Slide 10: Vielen Dank
add_content_slide(
    "Vielen Dank",
    "Kontaktinformationen: [Ihre Kontaktdaten]\nQ&A"
)

# Save the presentation in the current directory
file_path = "Netzwerkeinrichtung_Lord_of_the_Rings_Style.pptx"
prs.save(file_path)

file_path
